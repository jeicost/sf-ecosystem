"""
Rebuild Salsa Burgers PPTX:
1. Delete the 8 distorted pillar slides (indices 7-14)
2. Re-add them with correct aspect-ratio images
3. Add Month 1 campaign slides (overview + 4 weeks)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

pptx_path = '/Users/carlosjacoste/Desktop/Claude/clients/salsa-burgers-agency-deck.pptx'
assets_base = '/Users/carlosjacoste/Desktop/Claude/assets/briefing-assets/pillars'

# ── COLORS ──
BLACK  = RGBColor(0x1A,0x1A,0x1A)
BLACK2 = RGBColor(0x11,0x11,0x11)
DARK2  = RGBColor(0x26,0x26,0x26)
DARK3  = RGBColor(0x2A,0x2A,0x2A)
CREAM  = RGBColor(0xF5,0xF0,0xE8)
CREAM2 = RGBColor(0xCC,0xC8,0xC0)
CREAM3 = RGBColor(0x88,0x84,0x7C)
ORANGE = RGBColor(0xFF,0x45,0x00)

PILLAR_COLORS = [
    RGBColor(0xFF,0x45,0x00),
    RGBColor(0xFF,0x6A,0x00),
    RGBColor(0xFF,0x8C,0x00),
    RGBColor(0xFF,0xA5,0x00),
    RGBColor(0xFF,0xB7,0x32),
    RGBColor(0xCC,0x37,0x00),
    RGBColor(0x99,0x33,0x00),
    RGBColor(0x66,0x22,0x00),
]

PILLARS = [
    {
        "num":"01","name":"DRIVE CRAVING","color":PILLAR_COLORS[0],
        "phase":"Consideration",
        "desc":"Pure appetite content. Cinematic close-ups, glossy sauce pours, burger textures, drip moments. Dark background only. The visual goal is to make the viewer stop scrolling and order immediately.",
        "format":"Single pic · FB/IG · Size 1200×1500 (4:5)",
        "ex1_tag":"APRIL — Content 6","ex1_title":"NOT A BORING SAUCE, NOT A BORING BURGER",
        "ex1_body":"Slide flow: 1) Hero burger beauty shot 2) Sauce close-up/pour 3) Burger + glove ritual detail 4) Packaging + plated burger 5) Bite-ready/drip moment 6) Final hero with sauce on side.\nCaption: \"This is not the kind of burger you scroll past. Big burger energy. Bold sauce attitude.\"",
        "ex2_tag":"APRIL — Content 5","ex2_title":"THE MALA",
        "ex2_body":"Mala sauce + Wagyu close-up. \"Have you ever tried a Mala burger? 🔥 Rich, bold, perfectly balanced with that signature tingle.\"\nCTA: 20% OFF first order on Grab.",
        "images":["01_drive_craving/glove_first_bite.jpg","01_drive_craving/duo_combo_explicit.jpg"],
    },
    {
        "num":"02","name":"RITUAL & PACKAGING","color":PILLAR_COLORS[1],
        "phase":"Awareness · Consideration",
        "desc":"The full SALSA experience — metallic bag, gloves, finishing sauce, first bite. This pillar teaches new customers how to 'be' Salsa. Every swipe is a ritual step. Albums perform best.",
        "format":"Photo album (6 slides) · FB/IG · Size 1200×1500",
        "ex1_tag":"APRIL — Content 3 & 14","ex1_title":"SALSA CLUB RULES",
        "ex1_body":"HL: \"WELCOME TO THE SALSA CLUB — Glove up. Pick your salsa. Pour your way.\"\nSlide flow: 1) Welcome/Glove up 2) Pick your salsa 3) Make it yours/Pour 4) First bite 5) Pass to your crew 6) Join the ritual\nHashtags: #SalsaBurgers #SalsaRitual #SathornEats",
        "ex2_tag":"FORMAT NOTE","ex2_title":"WHY ALBUMS WIN FOR RITUAL",
        "ex2_body":"Albums stop the scroll with the first image. Each slide = one ritual step. The swipe = the actual experience. Saves rate on ritual content is 3-4× higher than single posts. Most memorable format.",
        "images":["02_ritual_packaging/hero_ritual_setup.jpg","02_ritual_packaging/glove_bite_dark.jpg"],
    },
    {
        "num":"03","name":"BRAND CULT","color":PILLAR_COLORS[2],
        "phase":"Awareness · Conversion",
        "desc":"Bold, irreverent, pop-culture brand content that makes Salsa feel like a movement, not a menu. References cinema, nostalgia, streetwear, arcade energy. Builds a following that identifies with the brand, not just the food.",
        "format":"Album (6 slides) · FB/IG · Highly shareable · Tag-your-crew CTA",
        "ex1_tag":"APRIL — Content 1 & 12","ex1_title":"THE BLOCKBUSTER",
        "ex1_body":"\"Some burgers are made to be eaten. This one was made to become part of the scene. Late nights. Neon lights. Arcade energy. Welcome to the blockbuster burger factory.\"\nCTA: \"Tag the crew and grab your popcorn.\"\n#SalsaBurgers #SalsaCult #NotJustABurger",
        "ex2_tag":"APRIL — Content 4 & 9","ex2_title":"STRANGER BURGERS",
        "ex2_body":"\"An Upside Down Sauce Experience.\" Streetwear attitude, nostalgia-coded. \"We took a wrong turn… and ended up somewhere strange. Good news? They have Salsa Burgers here too.\"\nAvailable on IG + TikTok simultaneously.",
        "images":["03_brand_cult/blockbuster_store.jpg","03_brand_cult/stranger_kids_gamezone.jpg"],
    },
    {
        "num":"04","name":"TRUST & AUTHENTICITY","color":PILLAR_COLORS[3],
        "phase":"Loyalty · Social Proof",
        "desc":"Real people, real delivery, real experience. Riders as brand ambassadors, customer moments, behind the scenes. Converts fence-sitters — answers 'can I trust the delivery?' with a definitive yes.",
        "format":"Single pic · FB/IG · Story reposts · UGC",
        "ex1_tag":"APRIL — Content 8","ex1_title":"TRUST YOUR DRIVER",
        "ex1_body":"HL: \"TRUST YOUR DRIVER. He knows the route. He knows the ritual.\"\n\"Some drivers deliver burgers. A Salsa driver delivers the moment. Fast hands. Hot bag. No boring drop-off.\"\n#TrustYourDriver #SalsaRitual",
        "ex2_tag":"ONGOING STRATEGY","ex2_title":"RIDERS AS BRAND ASSET",
        "ex2_body":"Riders are trained, given chill area, water, wifi. When they feel respected, they recommend. Their photo in front of the pickup zone = authentic content. Their positive attitude = brand extension at every doorstep.",
        "images":["04_trust_authenticity/burger_bangkok_temple.jpg","04_trust_authenticity/burger_cafe_moment.jpg"],
    },
    {
        "num":"05","name":"SALSA PHRASES","color":PILLAR_COLORS[4],
        "phase":"Conversion · Brand Cult",
        "desc":"Shareable, funny, locally-relevant phrases that people tag their crew in. No product required — just the brand voice. Designed to spread via Story shares. Local humour, Thai context, universal appetite.",
        "format":"Single pic · FB/IG/TikTok · No CTA needed — the share IS the action",
        "ex1_tag":"APRIL — Content 10","ex1_title":"SALSA PHRASES SERIES",
        "ex1_body":"Active phrases in rotation:\n· \"DIP IT LIKE YOU MEAN IT\"\n· \"DIP NOW THINK LATER\"\n· \"YOUR HANGOVER FAVOURITE IN THAILAND\"\nAlso run in Thai language for local reach. Replicate accounts that post funny shareable phrases.",
        "ex2_tag":"FORMAT RULE","ex2_title":"BOLD TYPE, DARK BG, ONE LINE",
        "ex2_body":"Maximum visual impact = minimum words. Anton or Impact font. Full bleed dark background. Orange accent. The phrase IS the post. No food photo needed. Highest share rate of any format we run.",
        "images":["05_salsa_phrases/parental_advisory_explicit_dipping.jpg"],
    },
    {
        "num":"06","name":"SALSA ICONS","color":PILLAR_COLORS[5],
        "phase":"Awareness · Cultural",
        "desc":"Monthly cultural icon series. A real person — musician, athlete, chef, icon — who embodies the Salsa spirit: bold, authentic, built differently. Creates a recurring format fans anticipate each month.",
        "format":"Photo album · FB/IG · Monthly cadence · Educate + entertain",
        "ex1_tag":"APRIL — Content 5 & 11","ex1_title":"JUAN LUIS GUERRA",
        "ex1_body":"\"This month's SALSA ICON: Juan Luis Guerra. Santo Domingo 🇩🇴 1957 — The man who turned bachata and merengue into global culture. Won +20 Latin Grammy Awards.\"\nHits: Burbujas de Amor · Ojalá que Llueva Café · La Bilirrubina",
        "ex2_tag":"SERIES CONCEPT","ex2_title":"MONTHLY ROTATION",
        "ex2_body":"Each icon connects to brand DNA: Latino energy + bold personality + cultural crossover. Not always musicians — could be chefs, athletes, filmmakers. The link is always spirit, not genre. Fans start guessing who's next.",
        "images":["06_salsa_icons/juan_luis_guerra.jpg"],
    },
    {
        "num":"07","name":"NEWS & PROMOTIONS","color":PILLAR_COLORS[6],
        "phase":"Conversions · Awareness",
        "desc":"Grand openings, product launches, limited-time promos, combo reveals. Conversion-first content. Clear CTA, Grab link, urgency language. Controlled frequency so it stays impactful — never more than 30% of output.",
        "format":"Single pic · FB/IG · Hard CTA → Grab link · Urgency always present",
        "ex1_tag":"APRIL — Content 2, 3, 4, 7, 8","ex1_title":"GRAND OPENING + 20% OFF",
        "ex1_body":"HL: \"SATHORN… A new burger just landed.\"\nCTA: 20% OFF first order → r.grab.com/o/6ru2AJv2\nCOMBOS: Hero Combo · Full Craving · Duo Combo · Full Indulgence\nCaption: \"More sauce. More reasons to dip, bite, repeat.\"",
        "ex2_tag":"PROMO RULE","ex2_title":"NO PERMANENT DISCOUNTS",
        "ex2_body":"SALSA1 and 20% OFF are launch mechanics. They expire after Month 1. After launch, full price only. Permanent discounts kill premium perception. Limited drops (Sauce of Month) replace price promotions.",
        "images":["07_news_promotions/burger_tower_20off_tiktok.jpg"],
    },
    {
        "num":"08","name":"SALSA ICONIC MOMENTS","color":PILLAR_COLORS[7],
        "phase":"Awareness · Cultural",
        "desc":"Legendary sports, music, and pop-culture moments reimagined with Salsa. Burger replaces the ball. Sauce bottle replaces the trophy. Bold, witty, memeable. 'Inspired by' — never exact celebrity likenesses.",
        "format":"Photo album (6 slides) · FB/IG · 'Which moment should we remix next?' CTA",
        "ex1_tag":"APRIL — Content 7","ex1_title":"SAUCE MOMENTS",
        "ex1_body":"Slide flow: 1) The Hand of Sauce 2) The Last Dunk 3) Historic sauce pour 4) Trophy-lift with burger box 5) Victory pose with salsa bottle 6) End card / hero burger / \"Join the ritual\"\nCaption: \"Some moments changed the game. Ours just made it saucier.\"",
        "ex2_tag":"CREATIVE RULE","ex2_title":"INSPIRED-BY, NOT EXACT",
        "ex2_body":"Use recognisable body pose/action but replace ball/trophy with burger or sauce. Brand-safe. Culturally alive. Legally clean. Fan CTA: \"Which moment should we remix next?\"",
        "images":["08_iconic_moments/jordan_just_do_it_with_salsa.jpg","08_iconic_moments/maradona_hand_of_god.jpg"],
    },
]

# ── Month 1 Campaign Data ──
WEEK_THEME_COLOR = RGBColor(0xFF,0x45,0x00)
GREEN  = RGBColor(0x2E,0xCC,0x71)
BLUE   = RGBColor(0x30,0x9A,0xFF)
PURPLE = RGBColor(0xA0,0x55,0xFF)

MONTH1_WEEKS = [
    {
        "num":"01","dates":"MAY 5–11","theme":"THE RITUAL LANDS",
        "goal":"Maximize awareness in Sathorn · Establish the ritual in culture",
        "kpis":["300 IG followers","10K impressions","30–50 Grab orders","5 Grab reviews"],
        "content":[
            {"format":"REEL","day":"Mon 5","pillar":"P1 DRIVE CRAVING","color":PILLAR_COLORS[0],
             "title":"NOT A BORING BURGER","platform":"IG · FB",
             "desc":"15–30s dark studio. Sauce pour close-up. Bold text overlay. First bite reveal.\nCaption: \"This is not the kind of burger you scroll past. Dark. Bold. 18 sauces. 100% Wagyu.\""},
            {"format":"CAROUSEL","day":"Wed 7","pillar":"P2 RITUAL","color":PILLAR_COLORS[1],
             "title":"SALSA CLUB RULES v2","platform":"IG · FB",
             "desc":"6 slides: Glove up / Pick your salsa / Pour with confidence / First bite / Pass to crew / Join the ritual.\n\"WELCOME TO THE SALSA CLUB.\""},
            {"format":"STORY","day":"Tue–Thu","pillar":"P2 RITUAL","color":PILLAR_COLORS[1],
             "title":"Kitchen Reveal + Ritual Tutorial","platform":"IG Stories",
             "desc":"5-frame series: Kitchen setup → Sauce lineup (18 salsas) → Gloves unbox → Box reveal → Grab CTA.\nInteractive poll: \"Which sauce are you?\""},
            {"format":"TIKTOK","day":"Sat 10","pillar":"P2 RITUAL","color":PILLAR_COLORS[1],
             "title":"SALSA RITUAL ASMR","platform":"TikTok",
             "desc":"Satisfying ritual video: glove unboxing + sauce pour + first bite. Trending audio overlay.\nGoal: save rate + shares. Hook: first 2 seconds = gloves snapping on."},
            {"format":"PAID","day":"All week","pillar":"META + GRAB","color":RGBColor(0x22,0x22,0x22),
             "title":"Awareness Push — The Ritual","platform":"Meta · Grab",
             "desc":"Meta: 15s Ritual video. Target: expats Sathorn 25–42, interests: food delivery, premium dining.\nGrab: Featured placement + Search. Banner: \"20% OFF. First order. Try the ritual.\""},
        ]
    },
    {
        "num":"02","dates":"MAY 12–18","theme":"BUILD THE CULT",
        "goal":"Develop brand personality · Seed first creators · Build social proof",
        "kpis":["500 IG followers","20K impressions","50–70 Grab orders","10 Grab reviews · 3 creator posts live"],
        "content":[
            {"format":"REEL","day":"Mon 12","pillar":"P3 BRAND CULT","color":PILLAR_COLORS[2],
             "title":"STRANGER BURGERS Ep.1","platform":"IG · FB · TikTok",
             "desc":"Cinematic teaser. Arcade nostalgia aesthetic. \"An Upside Down Sauce Experience.\"\n\"We went somewhere strange... and found Salsa Burgers.\"\n#StrangerBurgers #NotJustABurger"},
            {"format":"STATIC","day":"Thu 15","pillar":"P4 TRUST","color":PILLAR_COLORS[3],
             "title":"TRUST YOUR DRIVER","platform":"IG · FB",
             "desc":"Rider portrait + branded bag. \"He knows the route. He knows the ritual. Fast hands. Hot bag. No boring drop-off.\"\n#TrustYourDriver #SalsaRitual"},
            {"format":"STORY","day":"Tue–Fri","pillar":"P5 PHRASES","color":PILLAR_COLORS[4],
             "title":"Which Salsa Are You? Quiz","platform":"IG Stories",
             "desc":"Interactive 5-flavor personality quiz: MALA=bold · Honey Mustard=classic · Chipotle=adventurous.\nFri: Creator Kit reveal — first Tier A creator content drops."},
            {"format":"TIKTOK","day":"Sat 17","pillar":"P5 PHRASES","color":PILLAR_COLORS[4],
             "title":"Salsa Personality Challenge","platform":"TikTok",
             "desc":"\"Which salsa personality are you?\" duet/stitch challenge setup.\nHook: \"Tag someone who is MALA energy 🔥\"\nGoal: user-generated duets, organic reach amplification."},
            {"format":"PAID","day":"All week","pillar":"META + TIKTOK","color":RGBColor(0x22,0x22,0x22),
             "title":"Consideration Layer","platform":"Meta · TikTok",
             "desc":"Meta: \"Which Flavor Are You?\" carousel ad targeting engaged audience from Week 1.\nTikTok Reach: Boost Stranger Burgers reel. Target: Thai foodies + expat Bangkok community."},
        ]
    },
    {
        "num":"03","dates":"MAY 19–25","theme":"CULTURE + PHRASES",
        "goal":"Drive cultural relevance · Phrase virality · Creator amplification",
        "kpis":["800 IG followers","35K impressions","70–90 Grab orders","20+ Grab reviews · 8 creator posts"],
        "content":[
            {"format":"ALBUM","day":"Mon 19","pillar":"P6 ICONS","color":PILLAR_COLORS[5],
             "title":"SALSA ICON: May Edition","platform":"IG · FB",
             "desc":"May icon: a bold cultural figure connecting Latino energy + Bangkok relevance (TBD with team).\n\"Bold. Built differently. A true SALSA ICON.\"\nAlbum: 6 slides bio + career highlights + connection to Salsa spirit."},
            {"format":"STATIC","day":"Wed 21","pillar":"P5 PHRASES","color":PILLAR_COLORS[4],
             "title":"DIP IT LIKE YOU MEAN IT","platform":"IG · FB · TikTok",
             "desc":"Full bleed dark background. Orange Anton font. No food. Just the phrase.\nThai version Thu 22: localized for urban Thai audience.\nGoal: maximum Story shares and saves."},
            {"format":"STORY","day":"Tue–Sat","pillar":"P4 TRUST","color":PILLAR_COLORS[3],
             "title":"Creator Content Drops + Review Push","platform":"IG Stories",
             "desc":"Repost creator Ritual Kit reactions and unboxings as Stories.\nThu: \"Rate your ritual\" — Guide followers to Grab review. Show screenshot examples.\n\"Your review = helping us level up.\""},
            {"format":"TIKTOK","day":"Fri 23","pillar":"P2 RITUAL","color":PILLAR_COLORS[1],
             "title":"Creator Collab: Ritual POV","platform":"TikTok · IG Reels",
             "desc":"Tier A food creator films authentic Ritual Kit unboxing + first bite POV.\nTrending audio. Raw and genuine — no over-production.\nGoal: saves rate + profile visits → Grab conversion."},
            {"format":"PAID","day":"All week","pillar":"TIKTOK + META","color":RGBColor(0x22,0x22,0x22),
             "title":"Follows + Retargeting","platform":"TikTok · Meta",
             "desc":"TikTok: Follows campaign. Creative: creator ritual video. Target: Thai foodies + expats.\nMeta Retargeting: \"Still thinking about it?\" targeting people who saw Week 1 video but didn't convert.\nCopy: \"20% OFF ends soon.\""},
        ]
    },
    {
        "num":"04","dates":"MAY 26–31","theme":"CONVERT + CELEBRATE",
        "goal":"Final conversion push · Hit review milestone · Close Month 1 strong",
        "kpis":["1,000+ IG followers","100K+ impressions","120–150 total orders","30+ reviews · 4.5⭐"],
        "content":[
            {"format":"ALBUM","day":"Mon 26","pillar":"P8 ICONIC MOMENTS","color":PILLAR_COLORS[7],
             "title":"SAUCE MOMENTS Vol.1","platform":"IG · FB",
             "desc":"6-slide album: iconic sports/culture moments reimagined with Salsa.\nSlide flow: Hand of Sauce → The Last Dunk → Trophy-lift with burger box → Victory pose with salsa bottle → End card hero burger.\n\"Which moment should we remix next?\""},
            {"format":"STATIC","day":"Wed 28","pillar":"P7 PROMOS","color":PILLAR_COLORS[6],
             "title":"LAST DAYS — 20% OFF","platform":"IG · FB",
             "desc":"Urgency static: countdown energy. Bold copy.\n\"Order before the ritual goes full price.\"\nCTA: Grab link. Hard deadline. This promo never comes back — brand goes full price from June."},
            {"format":"TIKTOK","day":"Fri 30","pillar":"P3 BRAND CULT","color":PILLAR_COLORS[2],
             "title":"Month 1 Recap — Bangkok Dips","platform":"TikTok · IG Reels",
             "desc":"\"This is how Bangkok dips.\" Compilation reel: customer moments + creator content + ritual clips.\nEmotional closer. Trending audio. Shows the movement is real.\nEnd frame: \"Month 1 done. This is just the beginning.\""},
            {"format":"STORY","day":"All week","pillar":"P4 TRUST","color":PILLAR_COLORS[3],
             "title":"Review Push + Founder Moment","platform":"IG Stories",
             "desc":"Daily countdown: \"3 days... 2 days... promo ends.\" Screenshot review examples.\nSat 31: Authentic founder + team celebration. \"Month 1 ✅. We're just getting started.\"\nGoal: 30+ Grab reviews before month closes."},
            {"format":"PAID","day":"All week","pillar":"GRAB HEAVY + META","color":RGBColor(0x22,0x22,0x22),
             "title":"Final Conversion Sprint","platform":"Grab · Meta",
             "desc":"Grab: Max budget for final promo days. Featured + Search + Discount banner.\nMeta: Last-chance retargeting. \"20% OFF ends [date]. Order now.\"\nGoal: push cumulative orders to 150+ and close strong."},
        ]
    },
]


# ────────────────────────────────────────────────────────
# HELPERS
# ────────────────────────────────────────────────────────

def add_textbox(slide, left, top, width, height, text, font_size=10, bold=False,
                color=None, align=PP_ALIGN.LEFT, font_name="Arial", wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color if color else CREAM
    return txBox


def add_rect(slide, left, top, width, height, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w: shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_picture_fit(slide, img_path, left, top, max_w, max_h):
    """Add image maintaining aspect ratio, centered in available space."""
    if not os.path.exists(img_path):
        return None
    with Image.open(img_path) as img:
        iw, ih = img.size
    scale = min(max_w / iw, max_h / ih)
    w = int(iw * scale)
    h = int(ih * scale)
    ox = (max_w - w) // 2
    oy = (max_h - h) // 2
    return slide.shapes.add_picture(img_path, left + ox, top + oy, w, h)


def delete_slide(prs, idx):
    """Delete slide at index idx."""
    sldIdLst = prs.slides._sldIdLst
    elem = list(sldIdLst)[idx]
    rId = elem.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    sldIdLst.remove(elem)
    try:
        del prs.part._rels[rId]
    except Exception:
        pass


def add_pillar_slide(prs, p):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    W, H = prs.slide_width, prs.slide_height
    M = Inches(0.42)

    # BG
    add_rect(slide, 0, 0, W, H, fill=BLACK)
    # Accent bar
    add_rect(slide, 0, 0, Inches(0.055), H, fill=p["color"])

    # Section label
    add_textbox(slide, M, Inches(0.27), Inches(5), Inches(0.2),
                "07b — CONTENT PILLARS", 8, color=ORANGE)

    # Number
    add_textbox(slide, M, Inches(0.50), Inches(1.1), Inches(0.75),
                p["num"], 46, bold=True, color=p["color"], font_name="Arial Black")

    # Name
    add_textbox(slide, M+Inches(1.1), Inches(0.56), Inches(5.8), Inches(0.48),
                p["name"], 22, bold=True, color=CREAM, font_name="Arial Black")

    # Phase
    add_textbox(slide, M+Inches(1.1), Inches(1.0), Inches(4), Inches(0.22),
                p["phase"].upper(), 8, color=p["color"])

    # Divider
    add_rect(slide, M, Inches(1.26), Inches(8.6), Emu(16000), fill=p["color"])

    # Description
    add_textbox(slide, M, Inches(1.36), Inches(8.6), Inches(0.9),
                p["desc"], 10, color=CREAM2)

    # Format bar
    add_rect(slide, M, Inches(2.22), Inches(8.6), Inches(0.26), fill=DARK3)
    add_textbox(slide, M+Inches(0.1), Inches(2.24), Inches(8.4), Inches(0.22),
                "FORMAT  ·  " + p["format"], 8.5, color=CREAM2)

    # Example boxes
    ex_top = Inches(2.55)
    ex_h   = Inches(1.72)
    ex_w   = Inches(4.18)
    gap    = Inches(0.14)

    for i, (tag_k, ttl_k, body_k) in enumerate([
        ("ex1_tag","ex1_title","ex1_body"),
        ("ex2_tag","ex2_title","ex2_body"),
    ]):
        xl = M + i * (ex_w + gap)
        add_rect(slide, xl, ex_top, ex_w, ex_h, fill=DARK2)
        add_rect(slide, xl, ex_top, ex_w, Emu(14000), fill=p["color"])  # top accent
        add_textbox(slide, xl+Inches(0.1), ex_top+Inches(0.03), ex_w-Inches(0.18), Inches(0.18),
                    p[tag_k], 7.5, color=ORANGE)
        add_textbox(slide, xl+Inches(0.1), ex_top+Inches(0.21), ex_w-Inches(0.18), Inches(0.26),
                    p[ttl_k], 9.5, bold=True, color=CREAM, font_name="Arial Black")
        add_textbox(slide, xl+Inches(0.1), ex_top+Inches(0.48), ex_w-Inches(0.18), Inches(1.2),
                    p[body_k], 8.2, color=CREAM2)

    # Images — right column, properly fitted
    img_left = M + ex_w*2 + gap*2
    img_col_w = W - img_left - Inches(0.18)
    img_top  = Inches(1.36)
    img_tot_h = H - img_top - Inches(0.08)

    imgs = p["images"]
    if len(imgs) == 1:
        img_path = os.path.join(assets_base, imgs[0])
        add_picture_fit(slide, img_path, img_left, img_top, img_col_w, img_tot_h)
    elif len(imgs) >= 2:
        half = (img_tot_h - Inches(0.05)) // 2
        for i, rel in enumerate(imgs[:2]):
            img_path = os.path.join(assets_base, rel)
            y = img_top + i*(half + Inches(0.05))
            add_picture_fit(slide, img_path, img_left, y, img_col_w, half)

    # Footer
    add_textbox(slide, M, H-Inches(0.30), Inches(7), Inches(0.26),
                f"SALSA BURGERS  ·  PILLAR {p['num']}  ·  CONTENT STRATEGY 2026",
                7, color=CREAM3)
    return slide


def add_month1_overview(prs):
    """Month 1 overview slide"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    W, H = prs.slide_width, prs.slide_height
    M = Inches(0.5)

    add_rect(slide, 0, 0, W, H, fill=BLACK)
    add_rect(slide, 0, 0, Inches(0.055), H, fill=ORANGE)

    add_textbox(slide, M, Inches(0.27), Inches(6), Inches(0.2),
                "08b — MONTH 1 MASTER PLAN", 8, color=ORANGE)
    add_textbox(slide, M, Inches(0.50), Inches(9), Inches(0.75),
                "MAY 2026  ·  LAUNCH MONTH", 36, bold=True, color=CREAM, font_name="Arial Black")
    add_textbox(slide, M, Inches(1.18), Inches(9), Inches(0.22),
                "4 weeks · 4 themes · all formats · ฿32,000 paid media", 9, color=CREAM2)

    add_rect(slide, M, Inches(1.46), W-M*2, Emu(14000), fill=ORANGE)

    # KPI boxes
    kpi_data = [
        ("150–250","GRAB ORDERS"),("1,000+","IG FOLLOWERS"),
        ("30+","GRAB REVIEWS"),("4.5 ⭐","RATING TARGET"),
    ]
    bw = (W - M*2 - Inches(0.36)) / 4
    for i, (val, lbl) in enumerate(kpi_data):
        bx = M + i*(bw + Inches(0.12))
        add_rect(slide, bx, Inches(1.62), bw, Inches(0.92), fill=DARK2)
        add_textbox(slide, bx+Inches(0.1), Inches(1.68), bw-Inches(0.2), Inches(0.45),
                    val, 26, bold=True, color=ORANGE, font_name="Arial Black", align=PP_ALIGN.CENTER)
        add_textbox(slide, bx+Inches(0.1), Inches(2.10), bw-Inches(0.2), Inches(0.3),
                    lbl, 7.5, color=CREAM2, align=PP_ALIGN.CENTER)

    # 4-week cards
    themes = [
        ("01","MAY 5–11","THE RITUAL\nLANDS",PILLAR_COLORS[0]),
        ("02","MAY 12–18","BUILD THE\nCULT",PILLAR_COLORS[2]),
        ("03","MAY 19–25","CULTURE +\nPHRASES",PILLAR_COLORS[4]),
        ("04","MAY 26–31","CONVERT +\nCELEBRATE",PILLAR_COLORS[6]),
    ]
    ww = (W - M*2 - Inches(0.36)) / 4
    wt = Inches(2.70)
    wh = Inches(2.72)
    for i, (num, dates, theme, col) in enumerate(themes):
        wx = M + i*(ww + Inches(0.12))
        add_rect(slide, wx, wt, ww, wh, fill=DARK2)
        add_rect(slide, wx, wt, ww, Inches(0.06), fill=col)
        add_textbox(slide, wx+Inches(0.12), wt+Inches(0.12), ww-Inches(0.22), Inches(0.2),
                    f"WEEK {num}  ·  {dates}", 7, color=col)
        add_textbox(slide, wx+Inches(0.12), wt+Inches(0.34), ww-Inches(0.22), Inches(0.7),
                    theme, 17, bold=True, color=CREAM, font_name="Arial Black")
        # Format tags
        fmts = ["REEL","CAROUSEL / STATIC","STORIES","TIKTOK","PAID MEDIA"]
        for j, fmt in enumerate(fmts):
            fy = wt + Inches(1.08) + j*Inches(0.3)
            add_rect(slide, wx+Inches(0.12), fy, Inches(0.06), Inches(0.18), fill=col)
            add_textbox(slide, wx+Inches(0.22), fy, ww-Inches(0.34), Inches(0.22),
                        fmt, 7.5, color=CREAM2)

    add_textbox(slide, M, H-Inches(0.30), Inches(7), Inches(0.26),
                "SALSA BURGERS  ·  MONTH 1 MASTER PLAN  ·  NC GLOBAL ASSETS",
                7, color=CREAM3)
    return slide


def add_week_slide(prs, week):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    W, H = prs.slide_width, prs.slide_height
    M = Inches(0.42)

    add_rect(slide, 0, 0, W, H, fill=BLACK)
    col = PILLAR_COLORS[(int(week["num"])-1)*2 % 8]
    add_rect(slide, 0, 0, Inches(0.055), H, fill=col)

    # Header
    add_textbox(slide, M, Inches(0.2), W, Inches(0.2),
                f"08b — MONTH 1  ·  WEEK {week['num']}  ·  {week['dates']}", 8, color=ORANGE)
    add_textbox(slide, M, Inches(0.4), Inches(9), Inches(0.55),
                week["theme"], 30, bold=True, color=CREAM, font_name="Arial Black")
    add_textbox(slide, M, Inches(0.92), Inches(10), Inches(0.22),
                week["goal"], 9, color=CREAM2)

    # KPI pills
    kpi_x = M
    for kpi in week["kpis"]:
        kw = Inches(len(kpi)*0.075 + 0.3)
        add_rect(slide, kpi_x, Inches(1.16), kw, Inches(0.22), fill=col)
        add_textbox(slide, kpi_x+Inches(0.08), Inches(1.18), kw-Inches(0.1), Inches(0.18),
                    kpi, 7.5, bold=True, color=BLACK)
        kpi_x += kw + Inches(0.08)

    add_rect(slide, M, Inches(1.44), W-M*2, Emu(12000), fill=col)

    # Content rows
    row_h = Inches(0.94)
    rt = Inches(1.54)
    COL_WIDTHS = [Inches(1.0), Inches(0.72), Inches(1.72), Inches(2.68), Inches(1.75), Inches(3.9)]

    # Header row
    headers = ["FORMAT","DAY","PILLAR","TITLE + PLATFORM","DESCRIPTION",""]
    hx = M
    for hi, (hdr, hw) in enumerate(zip(headers, COL_WIDTHS)):
        if hdr:
            add_textbox(slide, hx, rt-Inches(0.24), hw, Inches(0.22),
                        hdr, 7, bold=True, color=CREAM3)
        hx += hw + Inches(0.04)

    for i, item in enumerate(week["content"]):
        y = rt + i * row_h
        bg_col = DARK2 if i%2==0 else RGBColor(0x20,0x20,0x20)
        add_rect(slide, M, y, W-M*2, row_h-Inches(0.04), fill=bg_col)
        # Left accent per pillar
        add_rect(slide, M, y, Inches(0.04), row_h-Inches(0.04), fill=item["color"])

        cx = M + Inches(0.06)
        vals = [
            item["format"], item["day"], item["pillar"],
            item["title"]+"\n"+item["platform"],
            item["desc"], ""
        ]
        for vi, (val, vw) in enumerate(zip(vals, COL_WIDTHS)):
            if vi == 0:  # FORMAT — bold pill style
                add_rect(slide, cx, y+Inches(0.32), vw-Inches(0.04), Inches(0.22), fill=item["color"])
                add_textbox(slide, cx, y+Inches(0.34), vw-Inches(0.04), Inches(0.18),
                            val, 7, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
            elif vi == 2:  # PILLAR
                add_textbox(slide, cx, y+Inches(0.08), vw, row_h-Inches(0.12),
                            val, 7, color=item["color"], bold=True)
            elif vi == 3:  # TITLE
                lines = val.split("\n")
                add_textbox(slide, cx, y+Inches(0.06), vw, Inches(0.32),
                            lines[0], 9, bold=True, color=CREAM, font_name="Arial Black")
                if len(lines) > 1:
                    add_textbox(slide, cx, y+Inches(0.38), vw, Inches(0.22),
                                lines[1], 7.5, color=item["color"])
            elif vi == 4:  # DESC
                add_textbox(slide, cx, y+Inches(0.04), vw, row_h-Inches(0.1),
                            val, 7.5, color=CREAM2)
            elif vi == 1:  # DAY
                add_textbox(slide, cx, y+Inches(0.08), vw, row_h-Inches(0.12),
                            val, 8, color=CREAM2)
            cx += vw + Inches(0.04)

    add_textbox(slide, M, H-Inches(0.28), Inches(8), Inches(0.24),
                f"SALSA BURGERS  ·  MAY 2026  ·  WEEK {week['num']}: {week['theme']}  ·  NC GLOBAL ASSETS",
                7, color=CREAM3)
    return slide


# ────────────────────────────────────────────────────────
# MAIN
# ────────────────────────────────────────────────────────
prs = Presentation(pptx_path)
print(f"Loaded PPTX: {len(prs.slides)} slides")

# Step 1: Delete slides 7–14 (the 8 old distorted pillar slides)
for _ in range(8):
    delete_slide(prs, 7)
print(f"After deletion: {len(prs.slides)} slides")

# Step 2: Add 8 fixed pillar slides (appended, then reordered)
for p in PILLARS:
    add_pillar_slide(prs, p)

# Step 3: Add Month 1 slides (appended)
add_month1_overview(prs)
for week in MONTH1_WEEKS:
    add_week_slide(prs, week)

# Step 4: Reorder — pillar slides (16-23) go after slide 6,
#         Month 1 slides (24-28) go after campaigns slide (originally slide 8, now at idx 7 after deletions)
# Current order after appending:
#   0-6:  original slides 1-7
#   7-15: original slides 8-16
#   16-23: 8 pillar slides
#   24-28: 5 Month1 slides
# Target:
#   0-6:  slides 1-7
#   7-14: 8 pillar slides
#   15:   slide 8 (campaigns) — original idx 7
#   16-20: 5 Month1 slides
#   21-23: slides 9-11  (original 9-16 minus the campaigns)
# Simpler target: 0-6, pillars 16-23, original 7-15, month1 24-28

sldIdLst = prs.slides._sldIdLst
all_ids = list(sldIdLst)
n = len(all_ids)
print(f"Total before reorder: {n}")

new_order = (
    list(range(7)) +         # original slides 1-7
    list(range(16, 24)) +    # 8 pillar slides
    [7] +                    # slide 8 (campaigns)
    list(range(24, 29)) +    # 5 Month1 slides
    list(range(8, 16))       # slides 9-16
)
print(f"New order ({len(new_order)} items): {new_order}")

for el in list(sldIdLst):
    sldIdLst.remove(el)
for idx in new_order:
    sldIdLst.append(all_ids[idx])

prs.save(pptx_path)
print(f"Saved. Final slide count: {len(prs.slides)}")
