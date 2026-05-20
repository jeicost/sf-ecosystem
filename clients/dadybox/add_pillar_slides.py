from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy, os

pptx_path = '/Users/carlosjacoste/Desktop/Claude/clients/salsa-burgers-agency-deck.pptx'
assets_base = '/Users/carlosjacoste/Desktop/Claude/assets/briefing-assets/pillars'
out_path = '/Users/carlosjacoste/Desktop/Claude/clients/salsa-burgers-agency-deck.pptx'

BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
CREAM  = RGBColor(0xF5, 0xF0, 0xE8)
CREAM2 = RGBColor(0xCC, 0xC8, 0xC0)
ORANGE = RGBColor(0xFF, 0x45, 0x00)

PILLARS = [
    {
        "num": "01", "name": "DRIVE CRAVING",
        "color": RGBColor(0xFF,0x45,0x00),
        "phase": "Consideration",
        "desc": "Pure appetite content. Cinematic close-ups, glossy sauce pours, burger textures, drip moments. Dark background only. The visual goal is to make the viewer stop scrolling and order immediately.",
        "format": "Single pic · FB/IG · Size 1200×1500 (4:5)",
        "ex1_tag": "APRIL — Content 6",
        "ex1_title": "NOT A BORING SAUCE, NOT A BORING BURGER",
        "ex1_body": "Slide flow: 1) Hero burger beauty shot 2) Sauce close-up/pour shot 3) Burger + glove ritual detail 4) Packaging + plated burger 5) Bite-ready angle / drip moment 6) Final hero with sauce on side.\nCaption: \"This is not the kind of burger you scroll past. Big burger energy. Bold sauce attitude.\"",
        "ex2_tag": "APRIL — Content 5",
        "ex2_title": "THE MALA",
        "ex2_body": "Mala sauce + Wagyu close-up. CTA: 20% OFF first order on Grab.",
        "images": ["01_drive_craving/glove_first_bite.jpg", "01_drive_craving/duo_combo_explicit.jpg"],
    },
    {
        "num": "02", "name": "RITUAL & PACKAGING",
        "color": RGBColor(0xFF,0x6A,0x00),
        "phase": "Awareness · Consideration",
        "desc": "The full SALSA experience — from metallic bag to gloves to finishing sauce to first bite. This pillar teaches new customers how to 'be' Salsa. Every swipe is a ritual step. Albums perform best.",
        "format": "Photo album (6 slides) · FB/IG · Size 1200×1500",
        "ex1_tag": "APRIL — Content 3 & 14",
        "ex1_title": "SALSA CLUB RULES",
        "ex1_body": "HL: \"WELCOME TO THE SALSA CLUB — Glove up. Pick your salsa. Pour your way.\"\nSlide flow: 1) Welcome/Glove up 2) Pick your salsa 3) Make it yours/Pour 4) First bite 5) Pass it to your crew 6) Join the ritual\nHashtags: #SalsaBurgers #SalsaRitual #SathornEats",
        "ex2_tag": "FORMAT NOTE",
        "ex2_title": "WHY ALBUMS WIN FOR RITUAL",
        "ex2_body": "Albums stop the scroll with the first image. Each slide = one ritual step. The swipe = the actual experience. Saves rate on ritual content is 3-4× higher than single posts.",
        "images": ["02_ritual_packaging/hero_ritual_setup.jpg", "02_ritual_packaging/glove_bite_dark.jpg"],
    },
    {
        "num": "03", "name": "BRAND CULT",
        "color": RGBColor(0xFF,0x8C,0x00),
        "phase": "Awareness · Conversion",
        "desc": "Bold, irreverent, pop-culture brand content that makes Salsa feel like a movement, not a menu. References cinema, nostalgia, streetwear, arcade energy. Builds a following that identifies with the brand, not just the food.",
        "format": "Album (6 slides) · FB/IG · Highly shareable · Tag-your-crew CTA",
        "ex1_tag": "APRIL — Content 1 & 12",
        "ex1_title": "THE BLOCKBUSTER",
        "ex1_body": "\"Some burgers are made to be eaten. This one was made to become part of the scene. Late nights. Neon lights. Arcade energy.\"\nCTA: \"Tag the crew and grab your popcorn.\"\n#SalsaBurgers #SalsaCult #NotJustABurger",
        "ex2_tag": "APRIL — Content 4 & 9",
        "ex2_title": "STRANGER BURGERS",
        "ex2_body": "\"An Upside Down Sauce Experience.\" Streetwear attitude, nostalgia-coded. Available on IG + TikTok simultaneously.\n#StrangerBurgers #NotJustABurger",
        "images": ["03_brand_cult/blockbuster_store.jpg", "03_brand_cult/stranger_kids_gamezone.jpg"],
    },
    {
        "num": "04", "name": "TRUST & AUTHENTICITY",
        "color": RGBColor(0xFF,0xA5,0x00),
        "phase": "Loyalty · Social Proof",
        "desc": "Real people, real delivery, real experience. Riders as brand ambassadors, customer moments, behind the scenes. This pillar converts fence-sitters — it answers 'but can I trust the delivery?' with a definitive yes.",
        "format": "Single pic · FB/IG · Story reposts · UGC",
        "ex1_tag": "APRIL — Content 8",
        "ex1_title": "TRUST YOUR DRIVER",
        "ex1_body": "HL: \"TRUST YOUR DRIVER. He knows the route. He knows the ritual.\"\n\"Some drivers deliver burgers. A Salsa driver delivers the moment. Fast hands. Hot bag. No boring drop-off.\"\n#TrustYourDriver #SalsaRitual",
        "ex2_tag": "ONGOING STRATEGY",
        "ex2_title": "RIDERS AS BRAND ASSET",
        "ex2_body": "Riders are trained, given chill area, water, wifi. When they feel respected, they recommend. Their photo in front of the pickup zone = authentic content.",
        "images": ["04_trust_authenticity/burger_bangkok_temple.jpg", "04_trust_authenticity/burger_cafe_moment.jpg"],
    },
    {
        "num": "05", "name": "SALSA PHRASES",
        "color": RGBColor(0xFF,0xB7,0x32),
        "phase": "Conversion · Brand Cult",
        "desc": "Shareable, funny, locally-relevant phrases that people tag their crew in. No product required — just the brand voice. Designed to spread via Story shares. Local humour, Thai context, universal appetite.",
        "format": "Single pic · FB/IG/TikTok · No CTA needed — share IS the action",
        "ex1_tag": "APRIL — Content 10",
        "ex1_title": "SALSA PHRASES SERIES",
        "ex1_body": "Active phrases in rotation:\n· \"DIP IT LIKE YOU MEAN IT\"\n· \"DIP NOW THINK LATER\"\n· \"YOUR HANGOVER FAVOURITE IN THAILAND\"\nAlso run in Thai language for local reach.",
        "ex2_tag": "FORMAT RULE",
        "ex2_title": "BOLD TYPE, DARK BG, ONE LINE",
        "ex2_body": "Maximum visual impact = minimum words. Anton or Impact font. Full bleed dark background. Orange accent. The phrase IS the post. No food photo needed. Highest share rate of any format.",
        "images": ["05_salsa_phrases/parental_advisory_explicit_dipping.jpg"],
    },
    {
        "num": "06", "name": "SALSA ICONS",
        "color": RGBColor(0xCC,0x37,0x00),
        "phase": "Awareness · Cultural",
        "desc": "Monthly cultural icon series. A real person — musician, athlete, chef, icon — who embodies the Salsa spirit: bold, authentic, built differently. Creates a recurring format fans anticipate each month.",
        "format": "Photo album · FB/IG · Monthly cadence · Educate + entertain",
        "ex1_tag": "APRIL — Content 5 & 11",
        "ex1_title": "JUAN LUIS GUERRA",
        "ex1_body": "\"This month's SALSA ICON: Juan Luis Guerra. Santo Domingo 🇩🇴 1957 — The man who turned bachata and merengue into global culture. Won +20 Latin Grammy Awards.\"\nHits: Burbujas de Amor · Ojalá que Llueva Café · La Bilirrubina",
        "ex2_tag": "SERIES CONCEPT",
        "ex2_title": "MONTHLY ROTATION",
        "ex2_body": "Each icon connects to brand DNA: Latino energy + bold personality + cultural crossover. Not always musicians — could be chefs, athletes, filmmakers. Fans start guessing who's next.",
        "images": ["06_salsa_icons/juan_luis_guerra.jpg"],
    },
    {
        "num": "07", "name": "NEWS & PROMOTIONS",
        "color": RGBColor(0x99,0x33,0x00),
        "phase": "Conversions · Awareness",
        "desc": "Grand openings, product launches, limited-time promos, combo reveals. Conversion-first content. Clear CTA, Grab link, urgency language. Kept to a controlled frequency so it stays impactful.",
        "format": "Single pic · FB/IG · Hard CTA → Grab link · Urgency always present",
        "ex1_tag": "APRIL — Content 2, 3, 4, 7, 8",
        "ex1_title": "GRAND OPENING + 20% OFF",
        "ex1_body": "HL: \"SATHORN… A new burger just landed.\"\nCTA: 20% OFF first order → r.grab.com/o/6ru2AJv2\nCOMBOS: Hero Combo · Full Craving · Duo Combo · Full Indulgence",
        "ex2_tag": "PROMO RULE",
        "ex2_title": "NO PERMANENT DISCOUNTS",
        "ex2_body": "SALSA1 and 20% OFF are launch mechanics. They expire. After Month 1, the brand sells at full price. Permanent discounts damage premium perception.",
        "images": ["07_news_promotions/burger_tower_20off_tiktok.jpg"],
    },
    {
        "num": "08", "name": "SALSA ICONIC MOMENTS",
        "color": RGBColor(0x66,0x22,0x00),
        "phase": "Awareness · Cultural",
        "desc": "Legendary sports, music, and pop-culture moments reimagined with Salsa. Burger replaces the ball. Sauce bottle replaces the trophy. Bold, witty, memeable. 'Inspired by' — never exact celebrity likenesses.",
        "format": "Photo album (6 slides) · FB/IG · 'Which moment should we remix next?' CTA",
        "ex1_tag": "APRIL — Content 7",
        "ex1_title": "SAUCE MOMENTS",
        "ex1_body": "Slide flow: 1) The Hand of Sauce 2) The Last Dunk 3) Historic sauce pour 4) Trophy-lift with burger box 5) Victory pose with salsa bottle 6) End card / hero burger\nCaption: \"Some moments changed the game. Ours just made it saucier.\"",
        "ex2_tag": "CREATIVE RULE",
        "ex2_title": "INSPIRED-BY, NOT EXACT",
        "ex2_body": "Use recognisable body pose/action but replace ball/trophy with burger or sauce. Brand-safe. Culturally alive. Legally clean. Fan engagement: 'Which moment should we remix next?'",
        "images": ["08_iconic_moments/jordan_just_do_it_with_salsa.jpg", "08_iconic_moments/maradona_hand_of_god.jpg"],
    },
]


def add_textbox(slide, left, top, width, height, text, font_size=10, bold=False,
                color=None, align=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = CREAM
    return txBox


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_pillar_slide(prs, pillar, insert_after_idx):
    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    W = prs.slide_width
    H = prs.slide_height
    M = Inches(0.4)  # margin

    # Background
    bg = add_rect(slide, 0, 0, W, H, fill_color=BLACK)

    # Left accent bar (pillar color)
    accent_bar = add_rect(slide, 0, 0, Inches(0.06), H, fill_color=pillar["color"])

    # Top label: "07b — CONTENT PILLARS"
    add_textbox(slide, M, Inches(0.28), Inches(5), Inches(0.22),
                "07b — CONTENT PILLARS",
                font_size=8, color=ORANGE, font_name="Arial")

    # Number (big)
    add_textbox(slide, M, Inches(0.52), Inches(1.2), Inches(0.8),
                pillar["num"],
                font_size=44, bold=True, color=pillar["color"], font_name="Arial Black")

    # Pillar name
    name_left = M + Inches(1.15)
    add_textbox(slide, name_left, Inches(0.58), Inches(5.5), Inches(0.5),
                pillar["name"],
                font_size=22, bold=True, color=CREAM, font_name="Arial Black")

    # Phase tag
    add_textbox(slide, name_left, Inches(1.0), Inches(4), Inches(0.24),
                pillar["phase"].upper(),
                font_size=8, color=pillar["color"], font_name="Arial")

    # Divider line
    line = add_rect(slide, M, Inches(1.28), Inches(8.6), Emu(18000), fill_color=pillar["color"])

    # Description
    add_textbox(slide, M, Inches(1.38), Inches(8.6), Inches(0.85),
                pillar["desc"],
                font_size=10.5, color=CREAM2, font_name="Arial")

    # Format bar
    fmt_box = add_rect(slide, M, Inches(2.18), Inches(8.6), Inches(0.28),
                       fill_color=RGBColor(0x2A,0x2A,0x2A))
    add_textbox(slide, M + Inches(0.1), Inches(2.2), Inches(8.4), Inches(0.26),
                "FORMAT: " + pillar["format"],
                font_size=9, color=CREAM2, font_name="Arial")

    # --- Examples ---
    ex_top = Inches(2.56)
    ex_h = Inches(1.65)
    ex_w = Inches(4.2)

    # Example 1
    add_rect(slide, M, ex_top, ex_w, ex_h,
             fill_color=RGBColor(0x26,0x26,0x26))
    add_rect(slide, M, ex_top, ex_w, Emu(18000),
             fill_color=pillar["color"])
    add_textbox(slide, M + Inches(0.12), ex_top + Inches(0.04), ex_w - Inches(0.2), Inches(0.18),
                pillar["ex1_tag"],
                font_size=7.5, color=ORANGE, font_name="Arial")
    add_textbox(slide, M + Inches(0.12), ex_top + Inches(0.22), ex_w - Inches(0.2), Inches(0.26),
                pillar["ex1_title"],
                font_size=10, bold=True, color=CREAM, font_name="Arial Black")
    add_textbox(slide, M + Inches(0.12), ex_top + Inches(0.50), ex_w - Inches(0.2), Inches(1.1),
                pillar["ex1_body"],
                font_size=8.5, color=CREAM2, font_name="Arial")

    # Example 2
    ex2_left = M + ex_w + Inches(0.15)
    add_rect(slide, ex2_left, ex_top, ex_w, ex_h,
             fill_color=RGBColor(0x26,0x26,0x26))
    add_rect(slide, ex2_left, ex_top, ex_w, Emu(18000),
             fill_color=pillar["color"])
    add_textbox(slide, ex2_left + Inches(0.12), ex_top + Inches(0.04), ex_w - Inches(0.2), Inches(0.18),
                pillar["ex2_tag"],
                font_size=7.5, color=ORANGE, font_name="Arial")
    add_textbox(slide, ex2_left + Inches(0.12), ex_top + Inches(0.22), ex_w - Inches(0.2), Inches(0.26),
                pillar["ex2_title"],
                font_size=10, bold=True, color=CREAM, font_name="Arial Black")
    add_textbox(slide, ex2_left + Inches(0.12), ex_top + Inches(0.50), ex_w - Inches(0.2), Inches(1.1),
                pillar["ex2_body"],
                font_size=8.5, color=CREAM2, font_name="Arial")

    # --- Images (right column) ---
    img_left = M + ex_w * 2 + Inches(0.3)
    img_col_w = W - img_left - Inches(0.2)
    img_top_start = Inches(1.38)
    img_total_h = H - img_top_start - Inches(0.1)

    imgs = pillar["images"]
    if len(imgs) == 1:
        img_path = os.path.join(assets_base, imgs[0])
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, img_left, img_top_start, img_col_w, img_total_h)
    elif len(imgs) >= 2:
        half_h = (img_total_h - Inches(0.06)) // 2
        for i, img_rel in enumerate(imgs[:2]):
            img_path = os.path.join(assets_base, img_rel)
            if os.path.exists(img_path):
                top = img_top_start + i * (half_h + Inches(0.06))
                slide.shapes.add_picture(img_path, img_left, top, img_col_w, half_h)

    # Bottom pill: slide label
    add_textbox(slide, M, H - Inches(0.32), Inches(6), Inches(0.28),
                f"SALSA BURGERS  ·  PILLAR {pillar['num']}  ·  CONTENT STRATEGY 2026",
                font_size=7.5, color=RGBColor(0x66,0x66,0x66), font_name="Arial")

    return slide


# ── Main ──
prs = Presentation(pptx_path)

# Insert 8 pillar slides after slide index 7 (slide 8 in 1-based = after "07 CONTENT" slide)
# We add them in reverse so each inserts at position 8
# Actually, add_slide appends — we'll reorder using XML
slides_xml = prs.slides._sldIdLst

# First, add all 8 slides (they'll be at the end)
new_slide_ids = []
for pillar in PILLARS:
    s = add_pillar_slide(prs, pillar, 7)
    # Get its rId
    for rel in prs.slides._sldIdLst:
        pass
    new_slide_ids.append(s)

print(f"Added {len(new_slide_ids)} pillar slides. Total: {len(prs.slides)}")

# Reorder: move slides 17-24 (0-indexed 16-23) to be after slide 7 (0-indexed 6)
# i.e. new order: 0,1,2,3,4,5,6, 16,17,18,19,20,21,22,23, 7,8,9,10,11,12,13,14,15
sldIdLst = prs.slides._sldIdLst
all_sldIds = list(sldIdLst)
n = len(all_sldIds)
# Original 16 slides: indices 0-15
# New 8 slides: indices 16-23
new_order = list(range(7)) + list(range(16, 24)) + list(range(7, 16))
print(f"Reordering {n} slides: {new_order}")

for el in list(sldIdLst):
    sldIdLst.remove(el)

for idx in new_order:
    sldIdLst.append(all_sldIds[idx])

prs.save(out_path)
print(f"Saved to {out_path}")
print(f"Final slide count: {len(prs.slides)}")
